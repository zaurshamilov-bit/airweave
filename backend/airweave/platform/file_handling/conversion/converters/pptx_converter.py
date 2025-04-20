"""PowerPoint to Markdown converter for PPTX files."""

from typing import Any, Union

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from airweave.core.logging import logger
from airweave.platform.file_handling.conversion._base import (
    DocumentConverter,
    DocumentConverterResult,
)


class PptxConverter(DocumentConverter):
    """Converts PPTX files to Markdown. Supports heading, tables and images with alt text."""

    async def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        """Convert a PPTX file to markdown.

        Args:
            local_path: Path to the PPTX file
            **kwargs: Additional arguments passed to converters

        Returns:
            DocumentConverterResult containing the markdown text
        """
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        try:
            presentation = pptx.Presentation(local_path)
            md_content = await self._process_presentation(presentation)
            return DocumentConverterResult(title=None, text_content=md_content.strip())
        except Exception as e:
            logger.error(f"Error converting PPTX file {local_path}: {str(e)}")
            return None

    async def _process_presentation(self, presentation: Any) -> str:
        """Process the entire presentation and convert to markdown.

        Args:
            presentation: PPTX presentation object

        Returns:
            Markdown content for the entire presentation
        """
        md_content = ""
        slide_num = 0

        for slide in presentation.slides:
            slide_num += 1
            md_content += await self._process_slide(slide, slide_num)

        return md_content

    async def _process_slide(self, slide: Any, slide_num: int) -> str:
        """Process a single slide and convert to markdown.

        Args:
            slide: PPTX slide object
            slide_num: Slide number

        Returns:
            Markdown content for this slide
        """
        md_content = f"\n\n<!-- Slide number: {slide_num} -->\n"

        # Process slide content (shapes)
        title = slide.shapes.title
        for shape in slide.shapes:
            md_content += await self._process_shape(shape, title)

        # Process slide notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            md_content += "\n\n### Notes:\n" + slide.notes_slide.notes_text_frame.text

        return md_content

    async def _process_shape(self, shape: Any, title: Any) -> str:
        """Process a single shape and convert to markdown.

        Args:
            shape: PPTX shape object
            title: The slide's title shape for comparison

        Returns:
            Markdown content for this shape
        """
        # Pictures
        if self._is_picture(shape):
            alt_text = ""
            try:
                alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
            except Exception:
                pass
            filename = shape.name.replace(" ", "_") + ".jpg"
            return f"\n![{alt_text if alt_text else shape.name}]({filename})\n"

        # Tables
        if self._is_table(shape):
            return await self._convert_table_to_markdown(shape.table)

        # Text areas
        if shape.has_text_frame:
            if shape == title:
                return "# " + shape.text.strip() + "\n"
            else:
                return shape.text + "\n"

        return ""

    def _is_picture(self, shape: Any) -> bool:
        """Check if a shape is a picture.

        Args:
            shape: Shape object to check

        Returns:
            True if the shape is a picture, False otherwise
        """
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape: Any) -> bool:
        """Check if a shape is a table.

        Args:
            shape: Shape object to check

        Returns:
            True if the shape is a table, False otherwise
        """
        return shape.shape_type == MSO_SHAPE_TYPE.TABLE

    async def _convert_table_to_markdown(self, table: Any) -> str:
        """Convert a PowerPoint table to markdown format.

        Args:
            table: PowerPoint table object

        Returns:
            Markdown table representation
        """
        markdown_rows = []
        header = []

        # Process header row
        for cell in table.rows[0].cells:
            header.append(cell.text.strip())
        markdown_rows.append("| " + " | ".join(header) + " |")
        markdown_rows.append("|" + "|".join(["---" for _ in header]) + "|")

        # Process data rows
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            markdown_rows.append("| " + " | ".join(cells) + " |")

        return "\n" + "\n".join(markdown_rows) + "\n"
