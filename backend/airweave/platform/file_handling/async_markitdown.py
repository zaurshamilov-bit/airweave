"""Async implementation of MarkItDown for PPTX, XLSX and Image files."""

import base64
import json
import mimetypes
import os
import shutil
import subprocess
import tempfile
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Dict, Optional, Set, Union

import pandas as pd
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from airweave.core.config import settings
from airweave.core.logging import logger

# Initialize OpenAI client if API key is available
openai_client = None
if settings.OPENAI_API_KEY:
    from openai import AsyncOpenAI

    openai_client = AsyncOpenAI(api_key=settings.OPENAI_API_KEY)


class AsyncDocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = "", file_path: str = ""):
        """Initialize the AsyncDocumentConverterResult."""
        self.title: Union[str, None] = title
        self.text_content: str = text_content
        self.file_path: str = file_path


class AsyncDocumentConverter(ABC):
    """Abstract base class for all document converters."""

    @abstractmethod
    async def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, AsyncDocumentConverterResult]:
        """Convert a document to markdown text."""
        pass


class AsyncPptxConverter(AsyncDocumentConverter):
    """Converts PPTX files to Markdown. Supports heading, tables and images with alt text."""

    async def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, AsyncDocumentConverterResult]:
        """Convert a PPTX file to markdown.

        Args:
            local_path: Path to the PPTX file
            **kwargs: Additional arguments passed to converters

        Returns:
            AsyncDocumentConverterResult containing the markdown text
        """
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""
        presentation = pptx.Presentation(local_path)
        slide_num = 0

        for slide in presentation.slides:
            slide_num += 1
            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        pass
                    filename = shape.name.replace(" ", "_") + ".jpg"
                    md_content += f"\n![{alt_text if alt_text else shape.name}]({filename})\n"

                # Tables
                if self._is_table(shape):
                    md_content += await self._convert_table_to_markdown(shape.table)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.strip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                md_content += "\n\n### Notes:\n" + slide.notes_slide.notes_text_frame.text

        return AsyncDocumentConverterResult(title=None, text_content=md_content.strip())

    def _is_picture(self, shape: Any) -> bool:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape: Any) -> bool:
        return shape.shape_type == MSO_SHAPE_TYPE.TABLE

    async def _convert_table_to_markdown(self, table: Any) -> str:
        """Convert a PowerPoint table to markdown format."""
        markdown_rows = []
        header = []

        # Process header row
        for cell in table.rows[0].cells:
            header.append(cell.text.strip())
        markdown_rows.append("| " + " | ".join(header) + " |")
        markdown_rows.append("|" + "|".join(["---" for _ in header]) + "|")

        # Process data rows
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            markdown_rows.append("| " + " | ".join(cells) + " |")

        return "\n" + "\n".join(markdown_rows) + "\n"


class AsyncXlsxConverter(AsyncDocumentConverter):
    """Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table."""

    async def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, AsyncDocumentConverterResult]:
        """Convert an XLSX file to markdown.

        Args:
            local_path: Path to the XLSX file
            **kwargs: Additional arguments passed to converters

        Returns:
            AsyncDocumentConverterResult containing the markdown text
        """
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xlsx":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None)
        md_content = ""

        for sheet_name, df in sheets.items():
            md_content += f"## {sheet_name}\n"
            md_content += df.to_markdown(index=False) + "\n\n"

        return AsyncDocumentConverterResult(title=None, text_content=md_content.strip())


class AsyncImageConverter(AsyncDocumentConverter):
    """Converts images to markdown via extraction of metadata and optional LLM description."""

    async def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, AsyncDocumentConverterResult]:
        """Convert an image to markdown.

        Args:
            local_path: Path to the image file
            **kwargs: Additional arguments passed to converters

        Returns:
            AsyncDocumentConverterResult containing the markdown text
        """
        # Bail if not supported image type
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        md_content = ""

        # Add metadata if exiftool is available
        metadata = await self._get_metadata(local_path)
        if metadata:
            for field in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if field in metadata:
                    md_content += f"{field}: {metadata[field]}\n"

        # Try describing the image with LLM if client is provided
        llm_client = kwargs.get("llm_client")
        llm_model = kwargs.get("llm_model")
        if llm_client and llm_model:
            description = await self._get_llm_description(
                local_path, extension, llm_client, llm_model, prompt=kwargs.get("llm_prompt")
            )
            if description:
                md_content += f"\n# Description:\n{description.strip()}\n"

        return AsyncDocumentConverterResult(title=None, text_content=md_content.strip())

    async def _get_metadata(self, local_path: str) -> Optional[Dict[str, Any]]:
        """Get image metadata using exiftool if available."""
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None

        try:
            result = subprocess.run(
                [exiftool, "-json", local_path], capture_output=True, text=True
            ).stdout
            return json.loads(result)[0]
        except Exception:
            return None

    async def _get_llm_description(
        self, local_path: str, extension: str, client: Any, model: str, prompt: Optional[str] = None
    ) -> str:
        """Get image description from LLM if available."""
        if not prompt:
            prompt = "Write a detailed caption for this image."

        # Convert image to data URI
        with open(local_path, "rb") as image_file:
            content_type, _ = mimetypes.guess_type("dummy" + extension)
            if not content_type:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

        # Prepare messages for LLM
        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": data_uri}},
                ],
            }
        ]

        # Get response from LLM
        response = await client.chat.completions.create(model=model, messages=messages)
        return response.choices[0].message.content


class AsyncPdfConverter(AsyncDocumentConverter):
    """Converts PDF files to Markdown using PyPDF2 or pdfminer.six if available."""

    async def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, AsyncDocumentConverterResult]:
        """Convert a PDF file to markdown.

        Args:
            local_path: Path to the PDF file
            **kwargs: Additional arguments passed to converters

        Returns:
            AsyncDocumentConverterResult containing the markdown text
        """
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        md_content = ""
        title = None

        try:
            # Try using PyPDF2 first
            import PyPDF2

            with open(local_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)

                # Try to extract title from metadata
                if reader.metadata and hasattr(reader.metadata, "title"):
                    title = reader.metadata.title

                # Extract text from each page
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        md_content += f"\n\n## Page {i + 1}\n\n{page_text}\n"
        except ImportError:
            # Fall back to pdfminer.six if PyPDF2 is not available
            try:
                from pdfminer.high_level import extract_text

                text = extract_text(local_path)
                md_content = text
            except ImportError:
                # If neither library is available, use subprocess to call external tools
                try:
                    # Try pdftotext if available
                    if shutil.which("pdftotext"):
                        result = subprocess.run(
                            ["pdftotext", local_path, "-"], capture_output=True, text=True
                        )
                        md_content = result.stdout
                    else:
                        return AsyncDocumentConverterResult(
                            title=None,
                            text_content=(
                                "PDF conversion requires PyPDF2, pdfminer.six, or pdftotext."
                            ),
                        )
                except Exception as e:
                    logger.error(f"Error converting PDF with external tools: {str(e)}")
                    return None

        return AsyncDocumentConverterResult(title=title, text_content=md_content.strip())


class AsyncDocxConverter(AsyncDocumentConverter):
    """Converts DOCX files to Markdown using python-docx or mammoth."""

    async def _convert_with_python_docx(self, local_path: str) -> tuple[str, Optional[str]]:
        """Convert DOCX using python-docx library.

        Args:
            local_path: Path to the DOCX file

        Returns:
            Tuple of (markdown_content, title)
        """
        import docx

        md_content = ""
        title = None

        doc = docx.Document(local_path)

        # Try to extract title from core properties
        if doc.core_properties.title:
            title = doc.core_properties.title

        # Process paragraphs
        for para in doc.paragraphs:
            if para.text:
                # Check if it's a heading
                if para.style.name.startswith("Heading"):
                    level = (
                        int(para.style.name.replace("Heading", ""))
                        if para.style.name != "Heading"
                        else 1
                    )
                    md_content += f"\n{'#' * level} {para.text}\n"
                else:
                    md_content += f"{para.text}\n\n"

        # Process tables
        for table in doc.tables:
            md_table = []
            # Add header row
            header_row = []
            for cell in table.rows[0].cells:
                header_row.append(cell.text.strip())
            md_table.append("| " + " | ".join(header_row) + " |")
            md_table.append("|" + "|".join(["---" for _ in header_row]) + "|")

            # Add data rows
            for row in table.rows[1:]:
                cells = [cell.text.strip() for cell in row.cells]
                md_table.append("| " + " | ".join(cells) + " |")

            md_content += "\n" + "\n".join(md_table) + "\n\n"

        return md_content, title

    async def _convert_with_mammoth(self, local_path: str) -> str:
        """Convert DOCX using mammoth library.

        Args:
            local_path: Path to the DOCX file

        Returns:
            Markdown content
        """
        import mammoth

        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_markdown(docx_file)
            return result.value

    async def _convert_with_pandoc(self, local_path: str) -> str:
        """Convert DOCX using pandoc.

        Args:
            local_path: Path to the DOCX file

        Returns:
            Markdown content
        """
        if not shutil.which("pandoc"):
            raise RuntimeError("Pandoc is not installed")

        with tempfile.NamedTemporaryFile(suffix=".md") as temp_file:
            subprocess.run(["pandoc", local_path, "-o", temp_file.name], check=True)
            with open(temp_file.name, "r") as f:
                return f.read()

    async def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, AsyncDocumentConverterResult]:
        """Convert a DOCX file to markdown.

        Args:
            local_path: Path to the DOCX file
            **kwargs: Additional arguments passed to converters

        Returns:
            AsyncDocumentConverterResult containing the markdown text
        """
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        md_content = ""
        title = None

        try:
            # Try using python-docx first
            try:
                md_content, title = await self._convert_with_python_docx(local_path)
            except ImportError:
                # Fall back to mammoth if python-docx is not available
                try:
                    md_content = await self._convert_with_mammoth(local_path)
                except ImportError:
                    # If neither library is available, try using pandoc
                    try:
                        md_content = await self._convert_with_pandoc(local_path)
                    except Exception as e:
                        logger.error(f"Error converting DOCX with external tools: {str(e)}")
                        return AsyncDocumentConverterResult(
                            title=None,
                            text_content=(
                                "DOCX conversion requires python-docx, mammoth, or pandoc."
                            ),
                        )
        except Exception as e:
            logger.error(f"Error converting DOCX: {str(e)}")
            return None

        return AsyncDocumentConverterResult(title=title, text_content=md_content.strip())


class AsyncHtmlConverter(AsyncDocumentConverter):
    """Converts HTML files to Markdown using html2text."""

    async def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, AsyncDocumentConverterResult]:
        """Convert an HTML file to markdown.

        Args:
            local_path: Path to the HTML file
            **kwargs: Additional arguments passed to converters

        Returns:
            AsyncDocumentConverterResult containing the markdown text
        """
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        md_content = ""
        title = None

        try:
            # Try using html2text
            import html2text

            with open(local_path, "r", encoding="utf-8") as file:
                html_content = file.read()

            # Try to extract title
            import re

            title_match = re.search(
                r"<title>(.*?)</title>", html_content, re.IGNORECASE | re.DOTALL
            )
            if title_match:
                title = title_match.group(1).strip()

            # Convert to markdown
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = False
            h.ignore_tables = False
            md_content = h.handle(html_content)

        except ImportError:
            # Fall back to pandoc if html2text is not available
            try:
                if shutil.which("pandoc"):
                    with tempfile.NamedTemporaryFile(suffix=".md") as temp_file:
                        subprocess.run(
                            [
                                "pandoc",
                                local_path,
                                "-f",
                                "html",
                                "-t",
                                "markdown",
                                "-o",
                                temp_file.name,
                            ],
                            check=True,
                        )
                        with open(temp_file.name, "r") as f:
                            md_content = f.read()
                else:
                    # Simple fallback using BeautifulSoup
                    try:
                        from bs4 import BeautifulSoup

                        with open(local_path, "r", encoding="utf-8") as file:
                            soup = BeautifulSoup(file.read(), "html.parser")

                        # Try to extract title
                        if soup.title:
                            title = soup.title.string

                        # Extract text
                        md_content = soup.get_text()
                    except ImportError:
                        return AsyncDocumentConverterResult(
                            title=None,
                            text_content=(
                                "HTML conversion requires html2text, pandoc, or BeautifulSoup."
                            ),
                        )
            except Exception as e:
                logger.error(f"Error converting HTML with external tools: {str(e)}")
                return None

        return AsyncDocumentConverterResult(title=title, text_content=md_content.strip())


class AsyncTextConverter(AsyncDocumentConverter):
    """Converts plain text files to Markdown."""

    async def _convert_csv_to_markdown(self, local_path: str) -> str:
        """Convert CSV file to markdown table.

        Args:
            local_path: Path to the CSV file

        Returns:
            Markdown table content
        """
        import csv

        md_content = ""
        csv_data = []

        with open(local_path, "r", encoding="utf-8") as csv_file:
            csv_reader = csv.reader(csv_file)
            for row in csv_reader:
                csv_data.append(row)

        if csv_data:
            # Create header
            md_content += "| " + " | ".join(csv_data[0]) + " |\n"
            md_content += "|" + "|".join(["---" for _ in csv_data[0]]) + "|\n"

            # Add data rows
            for row in csv_data[1:]:
                md_content += "| " + " | ".join(row) + " |\n"

        return md_content

    async def _convert_json_to_markdown(self, content: str) -> str:
        """Convert JSON content to formatted markdown.

        Args:
            content: JSON content as string

        Returns:
            Formatted markdown content
        """
        import json

        parsed_json = json.loads(content)
        return "```json\n" + json.dumps(parsed_json, indent=2) + "\n```"

    async def _convert_xml_to_markdown(self, content: str) -> str:
        """Convert XML content to formatted markdown.

        Args:
            content: XML content as string

        Returns:
            Formatted markdown content
        """
        import xml.dom.minidom

        parsed_xml = xml.dom.minidom.parseString(content)
        return "```xml\n" + parsed_xml.toprettyxml() + "\n```"

    async def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, AsyncDocumentConverterResult]:
        """Convert a text file to markdown.

        Args:
            local_path: Path to the text file
            **kwargs: Additional arguments passed to converters

        Returns:
            AsyncDocumentConverterResult containing the markdown text
        """
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".txt", ".csv", ".json", ".xml"]:
            return None

        md_content = ""
        title = None

        try:
            # For CSV files, convert to markdown table
            if extension.lower() == ".csv":
                try:
                    md_content = await self._convert_csv_to_markdown(local_path)
                except Exception as e:
                    logger.error(f"Error converting CSV to markdown table: {str(e)}")
                    # Fallback to raw content
                    with open(local_path, "r", encoding="utf-8") as file:
                        md_content = file.read()

            # For JSON files, pretty print
            elif extension.lower() == ".json":
                try:
                    with open(local_path, "r", encoding="utf-8") as file:
                        content = file.read()
                    md_content = await self._convert_json_to_markdown(content)
                except Exception as e:
                    logger.error(f"Error formatting JSON: {str(e)}")
                    # Fallback to raw content
                    with open(local_path, "r", encoding="utf-8") as file:
                        content = file.read()
                    md_content = "```\n" + content + "\n```"

            # For XML files, pretty print
            elif extension.lower() == ".xml":
                try:
                    with open(local_path, "r", encoding="utf-8") as file:
                        content = file.read()
                    md_content = await self._convert_xml_to_markdown(content)
                except Exception as e:
                    logger.error(f"Error formatting XML: {str(e)}")
                    # Fallback to raw content
                    with open(local_path, "r", encoding="utf-8") as file:
                        content = file.read()
                    md_content = "```\n" + content + "\n```"

            # For plain text, just use the content
            else:
                with open(local_path, "r", encoding="utf-8") as file:
                    md_content = file.read()

        except Exception as e:
            logger.error(f"Error reading text file: {str(e)}")
            return None

        return AsyncDocumentConverterResult(title=title, text_content=md_content.strip())


class AsyncMarkItDown:
    """Async implementation of MarkItDown for PPTX, XLSX and Image files."""

    SUPPORTED_EXTENSIONS: Set[str] = {
        ".pdf",
        ".doc",
        ".docx",
        ".ppt",
        ".pptx",
        ".xls",
        ".xlsx",
        ".html",
        ".htm",
        ".txt",
        ".csv",
        ".json",
        ".xml",
        ".png",
        ".jpg",
        ".jpeg",
    }

    def __init__(self, llm_client: Optional[Any] = None, llm_model: Optional[str] = None):
        """Initialize the AsyncMarkItDown converter."""
        self._llm_client = llm_client
        self._llm_model = llm_model
        self._converters = [
            AsyncPptxConverter(),
            AsyncXlsxConverter(),
            AsyncImageConverter(),
            AsyncPdfConverter(),
            AsyncDocxConverter(),
            AsyncHtmlConverter(),
            AsyncTextConverter(),
        ]

    async def convert(self, file_path: str, **kwargs: Any) -> AsyncDocumentConverterResult:
        """Convert a file to markdown format.

        Args:
            file_path: Path to the file to convert
            **kwargs: Additional arguments passed to converters

        Returns:
            AsyncDocumentConverterResult containing the markdown text
        """
        if not self._is_supported(file_path):
            raise ValueError(f"Unsupported file type: {Path(file_path).suffix}")

        if not os.path.exists(file_path):
            logger.warning(f"File not found: {file_path}")
            raise FileNotFoundError(f"File not found: {file_path}")

        # Get file extension
        _, extension = os.path.splitext(file_path)
        if not extension:
            raise ValueError("File must have an extension")

        # Add default kwargs
        if "file_extension" not in kwargs:
            kwargs["file_extension"] = extension
        if "llm_client" not in kwargs and self._llm_client:
            kwargs["llm_client"] = self._llm_client
        if "llm_model" not in kwargs and self._llm_model:
            kwargs["llm_model"] = self._llm_model

        # Try each converter
        for converter in self._converters:
            try:
                result = await converter.convert(file_path, **kwargs)
                if result:
                    result.file_path = file_path  # Add file path to result
                    return result
            except Exception as e:
                logger.error(f"Error converting with {converter.__class__.__name__}: {str(e)}")
                continue

        raise ValueError(f"No converter found for file type: {extension}")

    def _is_supported(self, file_path: str) -> bool:
        """Check if the file extension is supported."""
        ext = Path(file_path).suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            logger.warning(f"Unsupported file extension: {ext} for file: {file_path}")
            return False
        return True


# Create singleton instance
markitdown = AsyncMarkItDown(llm_client=openai_client, llm_model="gpt-4o-mini")
